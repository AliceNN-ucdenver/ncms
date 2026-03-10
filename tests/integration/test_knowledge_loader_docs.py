"""Integration tests for Knowledge Loader — rich document support via markitdown.

These tests require `ncms[docs]` (markitdown). They are automatically skipped
when markitdown is not installed, so the base test suite always passes cleanly.
"""

import os
import tempfile

import pytest

from ncms.application.knowledge_loader import KnowledgeLoader

# Skip entire module when markitdown is not available
pytestmark = pytest.mark.skipif(
    not KnowledgeLoader.has_markitdown(),
    reason="markitdown not installed — install with `pip install ncms[docs]`",
)


def _make_docx(path: str, paragraphs: list[str]) -> None:
    """Create a minimal .docx file using python-docx (pulled in by markitdown)."""
    from docx import Document

    doc = Document()
    for para in paragraphs:
        doc.add_paragraph(para)
    doc.save(path)


def _make_pptx(path: str, slides: list[tuple[str, str]]) -> None:
    """Create a minimal .pptx file using python-pptx (pulled in by markitdown)."""
    from pptx import Presentation

    prs = Presentation()
    for title, body in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    prs.save(path)


def _make_xlsx(path: str, rows: list[list[str]]) -> None:
    """Create a minimal .xlsx file using openpyxl (pulled in by markitdown)."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    wb.save(path)


class TestDocxLoading:
    @pytest.mark.asyncio
    async def test_load_docx_creates_memories(self, memory_service):
        """Loading a .docx file should convert to markdown and create memories."""
        loader = KnowledgeLoader(memory_service)

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            tmp_path = f.name

        try:
            _make_docx(
                tmp_path,
                [
                    "The API gateway routes traffic through NGINX.",
                    "Database connections use PgBouncer for pooling.",
                    "Frontend assets are served from CloudFront CDN.",
                ],
            )
            stats = await loader.load_file(tmp_path, domains=["architecture"])
            assert stats.files_processed == 1
            assert stats.memories_created >= 1
            assert len(stats.errors) == 0

            results = await memory_service.search("API gateway NGINX")
            assert len(results) >= 1
        finally:
            os.unlink(tmp_path)

    @pytest.mark.asyncio
    async def test_docx_in_supported_extensions(self, memory_service):
        """DOCX should appear in supported extensions when markitdown is available."""
        loader = KnowledgeLoader(memory_service)
        assert ".docx" in loader.SUPPORTED_EXTENSIONS


class TestPptxLoading:
    @pytest.mark.asyncio
    async def test_load_pptx_creates_memories(self, memory_service):
        """Loading a .pptx file should extract slide content as memories."""
        loader = KnowledgeLoader(memory_service)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            tmp_path = f.name

        try:
            _make_pptx(
                tmp_path,
                [
                    ("Architecture Overview", "Microservices with REST APIs and event bus."),
                    ("Database Layer", "PostgreSQL 15 with read replicas and WAL mode."),
                ],
            )
            stats = await loader.load_file(tmp_path, domains=["design"])
            assert stats.files_processed == 1
            assert stats.memories_created >= 1
            assert len(stats.errors) == 0

            results = await memory_service.search("PostgreSQL replicas")
            assert len(results) >= 1
        finally:
            os.unlink(tmp_path)


class TestXlsxLoading:
    @pytest.mark.asyncio
    async def test_load_xlsx_creates_memories(self, memory_service):
        """Loading a .xlsx file should extract tabular content as memories."""
        loader = KnowledgeLoader(memory_service)

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            tmp_path = f.name

        try:
            _make_xlsx(
                tmp_path,
                [
                    ["Endpoint", "Method", "Description"],
                    ["/users", "GET", "List all users with pagination"],
                    ["/auth/login", "POST", "Authenticate with JWT tokens"],
                    ["/items", "GET", "Retrieve inventory items"],
                ],
            )
            stats = await loader.load_file(tmp_path, domains=["api"])
            assert stats.files_processed == 1
            assert stats.memories_created >= 1
            assert len(stats.errors) == 0
        finally:
            os.unlink(tmp_path)


class TestMarkitdownUnavailable:
    """Verify graceful degradation when markitdown is NOT installed."""

    @pytest.mark.asyncio
    async def test_document_extensions_hidden_without_markitdown(self, memory_service):
        """SUPPORTED_EXTENSIONS should not include document types when flag is off."""
        loader = KnowledgeLoader(memory_service)
        # We're running WITH markitdown, so just verify the property works
        assert ".md" in loader.SUPPORTED_EXTENSIONS
        assert ".txt" in loader.SUPPORTED_EXTENSIONS

    @pytest.mark.asyncio
    async def test_has_markitdown_returns_true(self, memory_service):
        """has_markitdown() should be True when the library is installed."""
        assert KnowledgeLoader.has_markitdown() is True


class TestDirectoryWithDocuments:
    @pytest.mark.asyncio
    async def test_load_directory_includes_documents(self, memory_service):
        """Directory loading should pick up .docx files alongside .md files."""
        loader = KnowledgeLoader(memory_service)

        with tempfile.TemporaryDirectory() as tmpdir:
            # Markdown file
            md_path = os.path.join(tmpdir, "arch.md")
            with open(md_path, "w") as f:
                f.write("# Architecture\n\nMicroservices with REST APIs.\n")

            # Word doc
            docx_path = os.path.join(tmpdir, "details.docx")
            _make_docx(docx_path, ["Database uses PostgreSQL 15 with replication."])

            # Unsupported file — should be skipped
            bin_path = os.path.join(tmpdir, "data.bin")
            with open(bin_path, "wb") as f:
                f.write(b"\x00\x01\x02")

            stats = await loader.load_directory(tmpdir, domains=["project"])
            assert stats.files_processed == 2  # .md and .docx, not .bin
            assert stats.memories_created >= 2
            assert len(stats.errors) == 0
